import argparse
import random
import pptx
import odpslides.page
import odpslides.presentation


class Bullet():
    def __init__(self, text, level=0) -> None:
        self.text = text
        self.level = level

class Slide():
    def __init__(self, title, bullets) -> None:
        self.title = title
        self.bullets = bullets
    
    def add_bullet(self, bullet: Bullet):
        self.bullets.append(bullet)

class Presentation():
    def __init__(self, title, author="", slides=None) -> None:
        self.title = title
        self.author = author
        self.slides = slides if slides else []
    
    def add_slide(self, slide: Slide) -> None:
        self.slides.append(slide)

def lstrip_line(line: str):
    for c in ['#', '-', '*', '\t', ' ']:
        line = line.lstrip(c)
    return line.rstrip('\n')

def is_author(line: str):
    return line.lower().startswith('author:')

def get_author(line: str):
    return line.split(':')[1].lstrip(' ').rstrip('\n')

def get_intendation_level(line: str):
    level = 0
    line = line.expandtabs(4)
    while line.startswith(" " * 4):
        line = line.replace(" ", "", 4)
        level += 1
    return level

def parse_markdown(mdfn) -> Presentation:
    pres = None
    with open(mdfn, 'r') as mdf:
        title = mdf.readline()
        if not title or not title.startswith('#'):
            raise Exception("Markdown file should start with title headline")
        pres = Presentation(lstrip_line(title))
        for line in mdf:
            if line.startswith('#'):
                slide = Slide(lstrip_line(line), [])
                pres.add_slide(slide)
            elif is_author(line) and not pres.author:
                pres.author = get_author(line)
            else:
                level = get_intendation_level(line)
                line = line.lstrip('\t').lstrip(' ')
                if line.startswith('-') or line.startswith('*'):
                    slide.add_bullet(Bullet(lstrip_line(line), level))
    return pres

def pptx_add_title_slide(pres, title, author=""):
    title_slide_layout = pres.slide_layouts[0]
    slide = pres.slides.add_slide(title_slide_layout)
    t = slide.shapes.title
    t.text = title
    if author:
        subtitle = slide.placeholders[1]
        subtitle.text = author

def pptx_add_section_slide(pres, slide):
    section_slide_layout = pres.slide_layouts[2]
    pptx_slide = pres.slides.add_slide(section_slide_layout)
    shapes = pptx_slide.shapes
    title_shape = shapes.title
    title_shape.text = slide.title

def pptx_add_bullet_slide(pres, slide):
    bullet_slide_layout = pres.slide_layouts[1]
    pptx_slide = pres.slides.add_slide(bullet_slide_layout)
    shapes = pptx_slide.shapes
    title_shape = shapes.title
    title_shape.text = slide.title
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    # use text frame for first paragraph
    p = tf
    for b in slide.bullets:
        p.text = b.text
        p.level = b.level
        p = tf.add_paragraph()

def to_pptx(presentation: Presentation, outfn=""):
    pres = pptx.Presentation()
    pptx_add_title_slide(pres, presentation.title, presentation.author)
    for s in presentation.slides:
        if len(s.bullets) == 0:
            pptx_add_section_slide(pres, s)
        else:
            pptx_add_bullet_slide(pres, s)
    if outfn:
        pres.save(outfn)
    return pres

def to_odp(presentation: Presentation, outfn=""):
    pres = odpslides.presentation.Presentation()
    pres.add_title_chart(title = presentation.title, subtitle=presentation.author)
    for s in presentation.slides:
        if len(s.bullets) == 0:
            pres.add_title_chart(title=s.title)
        else:
            outline = [b.level * "\t" + b.text for b in s.bullets]
            pres.add_titled_outline_chart(title=s.title, outline=outline)
    if outfn:
        pres.save(filename=outfn)
    return pres

def intend(text, level):
    return level * "\t" + text

def tex_escape(text):
	for c in ["$", "%", "&", "#", "{", "}"]:
		text = text.replace(c, "\\"+c)
	return text

def to_tex(presentation: Presentation, templatefn="template.tex", aspect="43", outfn=""):
    colors = ["mdred", "mdblue", "mdgreen", "mdyellow", "mdviolet"]
    template = ""
    titlecolor = random.choice(colors)
    c = titlecolor
    with open(templatefn, "r") as f:
        template = f.read()
    slides_tex = []
    for s in presentation.slides:
        c = random.choice([color for color in colors if color!=c])
        slides_tex.append("\\setbeamercolor{{background canvas}}{{bg={bgcolor}}}".format(bgcolor=c))
        slides_tex.append("\\begin{{frame}}{{{title}}}".format(title=tex_escape(s.title)))
        slides_tex.append("\t\\begin{itemize}")
        level = 0
        for b in s.bullets:
            while(level < b.level):
                slides_tex.append(intend("\\begin{itemize}", level+2))
                level += 1
            while(level > b.level):
                slides_tex.append(intend("\\end{itemize}", level+1))
                level -= 1
            slides_tex.append(intend("\\item {text}".format(text=tex_escape(b.text)), level+2))
        while(level >= 0):
            slides_tex.append(intend("\\end{itemize}", level+1))
            level -= 1
        slides_tex.append("\\end{frame}\n")
    slides_tex = "\n".join(slides_tex)
    template = template.format(slides=slides_tex, title=tex_escape(presentation.title),
							   author=tex_escape(presentation.author), aspect=aspect, titlecolor=titlecolor)
    if outfn:
        with open(outfn, "w") as f:
            f.write(template)
    return template

if __name__=="__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("input", type=str, help="input markdown file")
    parser.add_argument("-o", "--out", type=str, default="out", help="output file name")
    parser.add_argument("-f", "--format", choices=["pptx", "odp", "tex"],
        default=["tex"], nargs="+", help="formats")
    parser.add_argument("--aspect", choices=["4:3", "16:9", "16:10"], default="4:3", help="aspect ratio of the slides")
    args = parser.parse_args()

    pres = parse_markdown(args.input)

    if "pptx" in args.format:
        p = to_pptx(pres, args.out  + ".pptx")
    if "odp" in args.format:
        p = to_odp(pres, args.out + ".odp")
    if "tex" in args.format:
        to_tex(pres, outfn=args.out+".tex", aspect=args.aspect.replace(":", ""))
